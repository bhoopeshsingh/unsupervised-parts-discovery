#!/usr/bin/env python3
"""Build viva_presentation_aligned.pptx from viva_presentation.pptx (plan alignment)."""

from pptx import Presentation
from pptx.util import Inches, Pt
import shutil
import sys

SRC = "viva_presentation.pptx"
OUT = "viva_presentation_aligned.pptx"


def replace_all(prs, old: str, new: str) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and old in shape.text:
                shape.text = shape.text.replace(old, new)


def main() -> None:
    shutil.copy2(SRC, OUT)
    prs = Presentation(OUT)

    # --- Multi-slide string replacements (order: longest first) ---
    seq = [
        (
            "Oxford Pets / Stanford Cars / CUB-200  ·  14,000 images  ·  3 classes",
            "Oxford-IIIT Pet & STL-10  ·  12,000 images (4000×3)  ·  3 classes",
        ),
        (
            "Architecture designed, implemented, and ablated here · 14,000 images · Oxford Pets / Stanford Cars / CUB-200",
            "Architecture designed, implemented, and ablated here · 12,000 images · Oxford-IIIT Pet & STL-10",
        ),
        (
            "DINO ViT-S/8  ·  Oxford Pets  ·  CUB",
            "DINO ViT-S/8  ·  Oxford-IIIT Pet  ·  STL-10",
        ),
        ("CS4700 Dissertation Viva", "CS4700A Dissertation Viva"),
        ("CS4700 Dissertation", "CS4700A Dissertation"),
    ]
    for old, new in seq:
        replace_all(prs, old, new)

    replace_all(prs, "14k images\n3 classes", "12k images\n3 classes")
    # Concept head matches code: multiclass LogReg (OneClassSVM = raw baseline only in scripts)
    replace_all(prs, "OneClass\nSVM", "3-class\nLogReg")

    # --- Slide 1: official title + DIFE ---
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text
        if "Unsupervised Concept Discovery" in t:
            shape.text = t.replace(
                "Unsupervised Concept Discovery",
                "Learning Interpretable Feature Spaces",
            )
        if "for Interpretable Clinical AI" in shape.text:
            shape.text = shape.text.replace(
                "for Interpretable Clinical AI",
                "via Domain-Informed Encoders (DIFE)",
            )

    for shape in s1.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text.strip().startswith("Bhoopesh Nandan Singh") and "250142934" not in shape.text:
            shape.text = shape.text.replace(
                "Bhoopesh Nandan Singh",
                "Bhoopesh Nandan Singh  ·  ID 250142934",
            )
            break

    has_sup = any(
        shape.has_text_frame and "Roberto Alamino" in shape.text
        for shape in s1.shapes
    )
    if not has_sup:
        box = s1.shapes.add_textbox(Inches(0.5), Inches(6.82), Inches(9.2), Inches(0.4))
        box.text_frame.paragraphs[0].text = "Supervisor: Dr. Roberto Alamino"
        box.text_frame.paragraphs[0].font.size = Pt(14)

    # --- Slide 4: DIFE label ---
    replace_all(
        prs,
        "Visual (DINO) and Tabular/Lab (PSFTT)",
        "DIFE: Visual (DINO) and Tabular/Lab (PSFTT)",
    )

    # --- Slide 6: STL-10 caption ---
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text.strip().startswith("Top concepts:"):
            shape.text = "STL-10 example — " + shape.text
            break

    # --- Slide 7: Figure 1 note (subtitle line, not title) ---
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text.strip().startswith("Correct-class concepts dominate"):
            shape.text = (
                shape.text
                + "\n\n(STL-10 three-class setup; aligns with dissertation Figure 1.)"
            )
            break

    # --- Slide 9: Table 13 alignment — keep title (1), subtitle (2), stats line (5); remove grid ---
    s9 = prs.slides[8]
    shapes = list(s9.shapes)
    for i, shape in enumerate(shapes):
        if i in (1, 2, 5):
            continue
        try:
            shape._element.getparent().remove(shape._element)
        except Exception:
            pass

    for shape in s9.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text.strip().startswith("All 30 Discovered"):
            shape.text = "All 30 Discovered Lab Concepts — NHANES (thesis Table 13)"
        elif "Named by LLM" in shape.text:
            shape.text = (
                "Named by LLM from deviation profiles only — zero diagnosis labels at training. "
                "CBC, Biochem, Lipid: 10 GMM clusters each. Representative rows from Table 13:"
            )
        elif "30/30" in shape.text:
            shape.text = (
                "30/30 clusters named; >90% LLM acceptance (dissertation §6.1.4). "
                "Healthy Lipid cluster (healthy_population, n≈1,221) noted in thesis."
            )

    lines = [
        "CBC (examples)",
        "• cbc:5 — metabolic_syndrome (n=3,014)  ·  DM enrichment 1.53×",
        "• cbc:3 — pre_diabetes (n=1,350)  ·  1.50×",
        "• cbc:6 — hypercholesterolemia (n=1,011)  ·  0.30×",
        "",
        "Biochem (examples)",
        "• biochem:0 — metabolic_syndrome (n=3,578)  ·  1.46×",
        "• biochem:5 — postprandial_hyperglycemia (n=1,198)  ·  1.42×",
        "• biochem:8 — hyperlipidemia_pattern (n=766)  ·  0.39×",
        "",
        "Lipid (examples)",
        "• lipid:6 — mild_hyperglycemia (n=1,131) — max DM enrichment 1.59×",
        "• lipid:9 — metabolic_syndrome (n=2,740)  ·  1.54×",
        "• lipid:7 — hyperlipidemia_pattern (n=2,588)  ·  0.45×",
        "",
        "See dissertation Table 13 for all 30 cluster IDs and labels.",
    ]
    tb = s9.shapes.add_textbox(Inches(0.45), Inches(1.5), Inches(9.1), Inches(5.35))
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
    tf.word_wrap = True

    # --- Slide 10: image results + tabular + trade-off narrative ---
    s10 = prs.slides[9]
    for shape in s10.shapes:
        if not shape.has_text_frame:
            continue
        if "Classification Results" in shape.text:
            shape.text = (
                "Classification Results: Image vs Tabular Domains"
            )
        elif shape.text.strip().startswith("LogReg probe"):
            shape.text = (
                "Image (STL-10): DIFE concept scores vs raw L12 baseline — test accuracy 98.9%, "
                "macro F1 = 0.99, AUC-ROC 0.99 (thesis Table 11; per-class bird / car / cat balanced). "
                "No accuracy–interpretability trade-off observed on this task.\n\n"
                "Tabular (NHANES): LogReg on 30 concept scores — no diagnosis labels used during concept learning."
            )
        elif shape.text.strip().startswith("Interpretation"):
            shape.text = (
                "Interpretation\n"
                "• Image: strong headroom — interpretable 30-dim concepts beat raw 384-dim features.\n"
                "• Tabular: 50.9% = 1.53× chance — real signal but classifier collapses toward 'normal'; "
                "GMM optimises cluster quality, not decision boundaries.\n"
                "• Trade-off: mainly in the clinical tabular regime (minority-class recall); "
                "the Abstract’s ‘no trade-off’ claim matches the image-domain result vs raw baseline."
            )
        elif "Trade-off is explicit: lower accuracy" in shape.text:
            shape.text = ""

    # --- Slide 5: masking footnote (θ grid search on Oxford-IIIT Pet) ---
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if not shape.has_text_frame:
            continue
        if "Attention-guided FG masking" in shape.text:
            shape.text = (
                shape.text
                + "\n\n(FG mask θ=0.75 from grid search on Oxford-IIIT Pet — thesis §4.2.3.)"
            )
            break

    # --- Slide 13: future work — CUB / Stanford Cars ---
    replace_all(
        prs,
        "Extend to MIMIC-IV for physician-confirmed diagnosis codes at scale",
        "Extend to MIMIC-IV for physician-confirmed diagnosis codes at scale\n"
        "Multi-dataset generalisation: CUB-200 & Stanford Cars (extension beyond Oxford-IIIT Pet / STL-10)",
    )

    # --- Slide 14: thank you title + metrics ---
    s14 = prs.slides[13]
    for shape in s14.shapes:
        if not shape.has_text_frame:
            continue
        if "Unsupervised Concept Discovery" in shape.text:
            shape.text = shape.text.replace(
                "Unsupervised Concept Discovery for Interpretable Clinical AI",
                "Learning Interpretable Feature Spaces (DIFE)",
            )

    # Add small metrics line if we can find email shape
    for shape in s14.shapes:
        if not shape.has_text_frame:
            continue
        if "bhoopesh.singh@gmail.com" in shape.text:
            shape.text = (
                shape.text
                + "\n\nImage: 98.9% acc · Tabular: 50.9% · NHANES n=18,673"
            )
            break

    prs.save(OUT)
    print("Wrote", OUT)


if __name__ == "__main__":
    import os

    os.chdir(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    main()
